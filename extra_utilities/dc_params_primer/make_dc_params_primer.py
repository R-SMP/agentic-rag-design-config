# -*- coding: utf-8 -*-
"""Build the design-configurator parameter primer.

Emits BOTH from one layout description:
  * dc_params_primer.pptx  - one slide of native, editable PowerPoint shapes
  * dc_params_primer.png   - the same drawing, for pasting into a prompt

    python make_dc_params_primer.py            # both
    python make_dc_params_primer.py --dpi 84   # smaller PNG

Panel A : top view  - ring, mid-wall diameter, one blade, and the shaded
                      green band the middle section may sit anywhere in.
Panel B : 3x3 grid  - camber across, high-point down.  Sections are drawn
                      chord-horizontal so the nine compare directly; the
                      centre cell names the chord line, the camber line and
                      the angle of attack (measured from the horizontal).

The airfoil math is a port of web/feg/naca.js (the same code path as
tools/generate_mesh/ring_height.py and tools/render_blade_sections), so the
shapes are the ones the DC actually builds.  Sections are drawn TE-left /
LE-right to match the render_blade_sections output the agents already see.

All geometry is in INCHES.
"""
from __future__ import division

import math
import os
import sys

from PIL import Image, ImageDraw, ImageFont

# ---------------------------------------------------------------------------
# 1.  Airfoil math  (port of web/feg/naca.js)
# ---------------------------------------------------------------------------


def _sym_profile(count_i, thickness_pct):
    n = count_i + 1
    cosv = [math.cos(i / (n - 1.0)) for i in range(n)]
    cs, ce = cosv[0], cosv[n - 1]
    remap = [(v - cs) / (ce - cs) for v in cosv]
    t_n = thickness_pct * 0.01
    z = []
    for y in remap:
        sy = 0.0 if y <= 0 else math.sqrt(y)
        z.append((t_n / 0.2) * (0.2969 * sy - 0.1260 * y - 0.3516 * y * y
                                + 0.2843 * y ** 3 - 0.1015 * y ** 4))
    pts = [(remap[i], -z[i]) for i in range(n - 1, -1, -1)]
    pts += [(remap[i], z[i]) for i in range(1, n)]
    return pts


def camber_curve(high_point_dec, camber_pct):
    p = max(0.001, min(0.999, high_point_dec * 0.1))
    m = camber_pct * 0.01
    out = []
    for i in range(21):
        x = i / 20.0
        if x <= p:
            z = (m / (p * p)) * (2 * p * x - x * x)
        else:
            z = (m / ((1 - p) ** 2)) * ((1 - 2 * p) + 2 * p * x - x * x)
        out.append((x, z))
    return out


def _morph(profile, camber):
    n = len(camber)
    res = []
    for py, pz in profile:
        f = py * (n - 1)
        seg = min(n - 2, max(0, int(math.floor(f))))
        u = min(1.0, max(0.0, f - seg))
        c0y, c0z = camber[seg]
        c1y, c1z = camber[seg + 1]
        cy = c0y + (c1y - c0y) * u
        cz = c0z + (c1z - c0z) * u
        ty, tz = c1y - c0y, c1z - c0z
        tl = math.hypot(ty, tz) or 1.0
        res.append((cy + pz * (-tz / tl), cz + pz * (ty / tl)))
    return res


def section(thickness, high_pt, camber, count_i=60):
    return _morph(_sym_profile(count_i, thickness), camber_curve(high_pt, camber))


# ---------------------------------------------------------------------------
# 2.  Display list
# ---------------------------------------------------------------------------

INK   = "1A1A1A"
GREY  = "7A7A7A"
FAINT = "C8C8C8"
BLUE  = "1F5FD0"      # inner section
GREEN = "1E8E4A"      # middle section / the middlePos band
CAMBER = "C00000"     # the camber line
OUTER  = "FF0000"     # outer section
ORANGE = "F5801F"     # ring mid-wall, and the diameter on it
BANDF = "EAF7EF"      # middlePos band fill
FOILL = "1F3A5F"
FOILF = "E8EEF5"
BLADE = "EFF2F6"
BLADL = "9AA4B0"
HUBF  = "EDEDED"      # hub fill — NEUTRAL: blue is reserved for the inner
                      # section, and the hub is not the inner section
WHITE = "FFFFFF"

# dash patterns: name -> (on, off) in inches, plus the PowerPoint preset
DASHES = {"short": (0.035, 0.026, "sysDash"),
          "dash":  (0.055, 0.036, "dash"),
          "long":  (0.090, 0.050, "lgDash")}

_FONT_CACHE = {}


def _font(pt, bold, dpi=96):
    key = (round(pt * dpi), bold)
    if key not in _FONT_CACHE:
        path = ("C:/Windows/Fonts/arialbd.ttf" if bold
                else "C:/Windows/Fonts/arial.ttf")
        _FONT_CACHE[key] = ImageFont.truetype(
            path, max(1, int(round(pt * dpi / 72.0))))
    return _FONT_CACHE[key]


def text_w(s, pt, bold=False):
    """Width of `s` in inches, measured with the real Arial metrics."""
    return _font(pt, bold, dpi=384).getlength(s) / 384.0


class Sheet(object):
    """Accumulates backend-neutral drawing primitives, in inches."""

    def __init__(self, w, h):
        self.w, self.h = w, h
        self.items = []

    def line(self, p0, p1, color=INK, pt=1.0, dash=None,
             arrow_start=False, arrow_end=False):
        self.items.append(("line", p0, p1, dict(
            color=color, pt=pt, dash=dash,
            arrow_start=arrow_start, arrow_end=arrow_end)))

    def circle(self, c, r, color=INK, pt=1.0, dash=None, fill=None,
               noline=False):
        self.items.append(("circle", c, r, dict(
            color=color, pt=pt, dash=dash, fill=fill, noline=noline)))

    def poly(self, pts, color=INK, pt=1.0, fill=None, close=True, dash=None):
        self.items.append(("poly", pts, None, dict(
            color=color, pt=pt, fill=fill, close=close, dash=dash)))

    def text(self, x, y, s, size=9.0, bold=False, color=INK, anchor="l"):
        """`y` is the TOP of the line; `anchor` is l / c / r on `x`."""
        self.items.append(("text", (x, y), s, dict(
            size=size, bold=bold, color=color, anchor=anchor)))


# ---------------------------------------------------------------------------
# 3.  The drawing
# ---------------------------------------------------------------------------

SLIDE_W, SLIDE_H = 11.4, 3.80

# panel A
A_X0 = 0.20
CX, CY = 2.70, 1.90            # shifted right: the detail needs the left column
R_MID = 1.33                      # impellerRadius, at the ring mid-wall
K = R_MID / 70.0                  # inches per mm
RING_T = 5.0                      # impellerThickness, mm
# TWO DIFFERENT RADII, and conflating them was a real defect in the first
# version of this drawing:
#   HUB_MM   the hub cylinder itself (web/feg/constants.js CONSTANTS.hub).
#   ROOT_MM  the radial station of the INNER BLADE SECTION -- 4.0 mm, from
#            inner_profile.cs (constants.js innerRadiusFixed).  It sits
#            INSIDE the hub, and it is the origin middlePos measures from:
#            profiles.js:19 is `radius = 4.0 + (impellerRadius - 4.0) * t`.
# Changing ROOT_MM here would put the drawing out of step with the geometry.
#
# HUB_MM is 8.0 by the owner's decision (2026-08-22), NOT the 8.28 in
# constants.js: that value is commented "interface.cs placeholder" and the
# hub is a cosmetic cylinder that no parameter depends on, so the round
# number is the one worth teaching.  ROOT_MM is the opposite case -- it IS
# load-bearing, so it tracks the code exactly and smoke_test_dc_primer
# asserts the two stay equal.
HUB_MM = 8.0
ROOT_MM = 4.0
B1_MM, B2_MM = ROOT_MM + 0.3 * 66, ROOT_MM + 0.7 * 66
DIVIDER = 4.90

# Enlarged detail of the centre, bottom-left of panel A.  At the top view's
# own scale (K = 0.019 in/mm) the two radii are 0.15 and 0.08 inches, far too
# close to label; the detail redraws them ~2.2x larger with both dimensioned.
DET_CX, DET_CY = 0.72, 2.75
DET_K = 0.34 / HUB_MM             # inches per mm inside the detail

# panel B
B_X0 = 5.05
GUT = 0.95
COL_W = [1.38, 2.36, 1.38]
ROW_H = [0.78, 1.28, 0.78]
GRID_TOP = 0.88
CAMS = [2.0, 5.0, 9.0]
HPS = [2, 5, 8]
THICK = 18.0
CHORD = 1.24
AOA = 15.0

T_TITLE, T_HEAD, T_LAB, T_ANN = 14.0, 10.5, 9.0, 9.0


def polar(r_mm, deg):
    a = math.radians(deg)
    return (CX + r_mm * K * math.cos(a), CY - r_mm * K * math.sin(a))


def build():
    s = Sheet(SLIDE_W, SLIDE_H)

    # ---------------- PANEL A -------------------------------------------
    s.text(A_X0, 0.12, "TOP VIEW", T_TITLE, True)
    s.text(A_X0 + text_w("TOP VIEW", T_TITLE, True) + 0.12, 0.19,
           u"–  propeller seen along its rotation axis", T_ANN, color=GREY)

    r_in = (70 - RING_T / 2) * K
    r_out = (70 + RING_T / 2) * K

    # the band the middle section may sit anywhere in
    s.circle((CX, CY), B2_MM * K, fill=BANDF, noline=True)
    s.circle((CX, CY), B1_MM * K, fill=WHITE, noline=True)
    s.circle((CX, CY), B2_MM * K, GREEN, 1.0, "dash")
    s.circle((CX, CY), B1_MM * K, GREEN, 1.0, "dash")

    # one blade, pointing up
    up, dn = [], []
    for i in range(41):
        t = i / 40.0
        r = ROOT_MM + (70 - ROOT_MM) * t
        c = 4.0 + 18.0 * (t ** 0.75)
        ha = math.degrees((c / 2.0) / r)
        up.append(polar(r, 90 + ha))
        dn.append(polar(r, 90 - ha))
    s.poly(up + dn[::-1], BLADL, 1.25, BLADE)

    # ring: the two wall faces, and the mid-wall the diameter is taken on
    s.circle((CX, CY), r_out, INK, 1.5)
    s.circle((CX, CY), r_in, INK, 1.5)
    s.circle((CX, CY), R_MID, ORANGE, 1.0, "dash")
    # hub (grey, solid) and the inner-section station inside it (blue, dashed)
    s.circle((CX, CY), HUB_MM * K, GREY, 1.25, fill=HUBF)
    s.circle((CX, CY), ROOT_MM * K, BLUE, 1.0, "short")

    s.line((CX - R_MID, CY), (CX + R_MID, CY), INK, 1.25,
           arrow_start=True, arrow_end=True)
    for sgn in (-1, 1):
        s.line((CX + sgn * R_MID, CY - 0.07), (CX + sgn * R_MID, CY + 0.07),
               ORANGE, 1.5)

    # how far out the middle section may sit, marked on the blade
    ya, yb = polar(B1_MM, 90)[1], polar(B2_MM, 90)[1]
    s.line((CX, ya), (CX, yb), GREEN, 1.75, arrow_start=True, arrow_end=True)

    # station callouts, colour-matched
    tipy = polar(70, 90)[1]
    halfc = 11.0 * K
    s.line((CX - halfc, tipy + 0.012), (CX + halfc, tipy + 0.012), OUTER, 0.75)
    s.line((CX + halfc - 0.01, tipy + 0.011), (3.52, 0.42), OUTER, 0.75)
    s.text(3.56, 0.36, "outer section", T_LAB, True, OUTER)
    s.text(3.56, 0.50, "at impellerRadius", T_LAB, color=OUTER)

    ymid = (ya + yb) / 2.0
    s.line((1.10, ymid), (CX - 0.15, ymid), GREEN, 0.75)
    s.text(1.06, ymid - 0.15, "middle section", T_LAB, True, GREEN, anchor="r")
    s.text(1.06, ymid - 0.01, u"middlePos 0.3–0.7", T_LAB, color=GREEN,
           anchor="r")

    # ---- enlarged detail: the hub is NOT the inner section ---------------
    # Both circles are ~1 mm across in the view above, far too close to
    # label there; drawn again here with each radius dimensioned, because
    # the two being DIFFERENT is the single fact this corner exists for.
    # A thin leader ties the detail back to the centre it magnifies.
    s.line((1.02, 2.62), (CX - 0.13, CY + 0.10), FAINT, 0.75)
    s.text(A_X0, DET_CY - 0.51, "centre, enlarged", T_ANN, color=GREY)
    s.circle((DET_CX, DET_CY), HUB_MM * DET_K, GREY, 1.25, fill=HUBF)
    s.circle((DET_CX, DET_CY), ROOT_MM * DET_K, BLUE, 1.25, "short")

    # Each radius gets its OWN arrow from the centre, at a different angle.
    # Both on one horizontal line read as a single 4-to-8 span instead of two
    # radii, which is the opposite of the point.
    def _spoke(r_mm, deg, color, label, dx, dy, anchor="l"):
        a = math.radians(deg)
        ex = DET_CX + r_mm * DET_K * math.cos(a)
        ey = DET_CY - r_mm * DET_K * math.sin(a)
        s.line((DET_CX, DET_CY), (ex, ey), color, 1.0, arrow_end=True)
        s.text(ex + dx, ey + dy, label, T_LAB, True, color, anchor=anchor)

    # 180 and -45: NOT opposite angles.  Two collinear spokes look like one
    # arrow spanning 4 to 8, which is the opposite of what this shows.
    _spoke(ROOT_MM, 180, BLUE, "4", -0.03, -0.13, anchor="r")
    _spoke(HUB_MM, -45, GREY, "8", 0.02, -0.06)

    # Kept SHORT: past x ~1.65 these rows run into the diameter caption.
    s.text(A_X0, 3.16, "hub  r = 8 mm", T_LAB, True, color=GREY)
    s.text(A_X0, 3.30, "inner section  r = 4 mm", T_LAB, True, color=BLUE)
    s.text(A_X0, 3.44, "— inside the hub", T_ANN, color=GREY)

    s.text(CX, 3.34, u"Ø = 2 × impellerRadius", T_HEAD, True,
           color=ORANGE, anchor="c")
    s.text(CX, 3.54, u"taken on the ring MID-WALL (orange dashes)", T_ANN,
           color=GREY, anchor="c")

    s.line((DIVIDER, 0.30), (DIVIDER, 3.55), FAINT, 0.75)

    # ---------------- PANEL B -------------------------------------------
    s.text(B_X0, 0.12, "BLADE SECTION SHAPE", T_TITLE, True)
    s.text(B_X0 + text_w("BLADE SECTION SHAPE", T_TITLE, True) + 0.12, 0.19,
           u"–  one blade section, seen edge-on", T_ANN, color=GREY)

    colx = [B_X0 + GUT]
    for w in COL_W:
        colx.append(colx[-1] + w)
    rowy = [GRID_TOP]
    for h in ROW_H:
        rowy.append(rowy[-1] + h)

    for i, c in enumerate(CAMS):
        s.text((colx[i] + colx[i + 1]) / 2.0, 0.58, "camber %g %%" % c,
               T_HEAD, True, anchor="c")
    for i, hp in enumerate(HPS):
        s.text(colx[0] - 0.10, (rowy[i] + rowy[i + 1]) / 2.0 - 0.08,
               "high-point %d" % hp, T_HEAD, True, anchor="r")

    for r, hp in enumerate(HPS):
        for c, cam in enumerate(CAMS):
            _cell(s, colx[c], colx[c + 1], rowy[r], rowy[r + 1], cam, hp,
                  row=r, annotate=(r == 1 and c == 1))

    return s


def _cell(s, x0, x1, y0, y1, camber, hp, row=0, annotate=False):
    """One grid cell.  Sections are chord-horizontal so the nine compare."""
    cx = (x0 + x1) / 2.0
    cy = y0 + 0.48 if row == 1 else (y0 + y1) / 2.0
    half = CHORD / 2.0

    if annotate:
        s.poly([(x0 + 0.03, y0 + 0.03), (x1 - 0.03, y0 + 0.03),
                (x1 - 0.03, y1 - 0.03), (x0 + 0.03, y1 - 0.03)], FAINT, 0.75)
        # the horizontal plane, hinged at the TE so the wedge opens clear
        a = math.radians(AOA)
        s.line((cx - half, cy),
               (cx - half + 1.62 * math.cos(a), cy + 1.62 * math.sin(a)),
               GREY, 1.0, "long")
        rr = 0.66
        s.poly([(cx - half + rr * math.cos(math.radians(t)),
                 cy + rr * math.sin(math.radians(t)))
                for t in range(0, int(AOA) + 1, 3)], GREY, 1.25, close=False)

    s.poly([(cx + (0.5 - p[0]) * CHORD, cy - p[1] * CHORD)
            for p in section(THICK, hp, camber)], FOILL, 1.25, FOILF)

    e_te = 0.20 if annotate else 0.07
    s.line((cx - half - e_te, cy), (cx + half + 0.07, cy), INK, 1.25, "dash")
    s.poly([(cx + (0.5 - p[0]) * CHORD, cy - p[1] * CHORD)
            for p in camber_curve(hp, camber)], CAMBER, 1.5, close=False,
           dash="short")
    if not annotate:
        return

    s.line((cx - 0.20, cy + 0.28), (cx - 0.05, cy + 0.11), GREY, 0.75)
    s.text(cx - 0.62, cy + 0.30, "angle of attack", T_ANN, color=GREY)
    s.text(cx + 0.86, cy + 0.52, "horizontal plane", T_ANN, color=GREY,
           anchor="r")

    s.text(cx - half - 0.24, cy - 0.27, "Chord line", T_ANN, True)
    crest = camber_curve(hp, camber)[min(20, hp * 2)]
    px, py = cx + (0.5 - crest[0]) * CHORD, cy - crest[1] * CHORD
    s.line((px, py - 0.03), (px, py - 0.19), CAMBER, 0.75)
    s.text(px, cy - 0.36, "Camber line", T_ANN, True, CAMBER, anchor="c")

    s.text(cx - half - 0.24, cy + 0.05, "TE", T_ANN, color=GREY, anchor="r")
    s.text(cx + half + 0.11, cy - 0.19, "LE", T_ANN, color=GREY)


# ---------------------------------------------------------------------------
# 4.  PNG backend
# ---------------------------------------------------------------------------


def render_png(sheet, path, dpi=96, ss=4):
    def rgb(h):
        return tuple(int(h[i:i + 2], 16) for i in (0, 2, 4))

    scale = dpi * ss
    W, H = int(sheet.w * scale), int(sheet.h * scale)
    img = Image.new("RGB", (W, H), (255, 255, 255))
    d = ImageDraw.Draw(img)

    def X(v):
        return v * scale

    def lw(pt):
        return max(1, int(round(pt / 72.0 * scale)))

    def seg(p0, p1, color, w, dash):
        if not dash:
            d.line([X(p0[0]), X(p0[1]), X(p1[0]), X(p1[1])],
                   fill=color, width=w)
            return
        on, off, _ = DASHES[dash]
        ln = math.hypot(p1[0] - p0[0], p1[1] - p0[1])
        if ln <= 0:
            return
        ux, uy = (p1[0] - p0[0]) / ln, (p1[1] - p0[1]) / ln
        t = 0.0
        while t < ln:
            t2 = min(ln, t + on)
            d.line([X(p0[0] + ux * t), X(p0[1] + uy * t),
                    X(p0[0] + ux * t2), X(p0[1] + uy * t2)],
                   fill=color, width=w)
            t = t2 + off

    def head(tip, ang, color, size=0.105):
        a1, a2 = ang + math.radians(166), ang - math.radians(166)
        d.polygon([(X(tip[0]), X(tip[1])),
                   (X(tip[0] + size * math.cos(a1)),
                    X(tip[1] + size * math.sin(a1))),
                   (X(tip[0] + size * math.cos(a2)),
                    X(tip[1] + size * math.sin(a2)))], fill=color)

    for kind, a, b, o in sheet.items:
        col = rgb(o.get("color", INK))
        if kind == "line":
            seg(a, b, col, lw(o["pt"]), o.get("dash"))
            ang = math.atan2(b[1] - a[1], b[0] - a[0])
            if o.get("arrow_end"):
                head(b, ang, col)
            if o.get("arrow_start"):
                head(a, ang + math.pi, col)
        elif kind == "circle":
            cx, cy, r = a[0], a[1], b
            box = [X(cx - r), X(cy - r), X(cx + r), X(cy + r)]
            if o.get("fill"):
                d.ellipse(box, fill=rgb(o["fill"]))
            if o.get("noline"):
                continue
            if o.get("dash"):
                for deg in range(0, 360, 12):
                    d.arc(box, deg, deg + 6, fill=col, width=lw(o["pt"]))
            else:
                d.ellipse(box, outline=col, width=lw(o["pt"]))
        elif kind == "poly":
            pts = [(X(p[0]), X(p[1])) for p in a]
            if o.get("fill"):
                d.polygon(pts, fill=rgb(o["fill"]))
            seq = a + [a[0]] if o.get("close", True) else a
            if o.get("dash"):
                for i in range(len(seq) - 1):
                    seg(seq[i], seq[i + 1], col, lw(o["pt"]), o["dash"])
            else:
                d.line([(X(p[0]), X(p[1])) for p in seq], fill=col,
                       width=lw(o["pt"]), joint="curve")
        else:
            f = _font(o["size"], o["bold"], dpi=scale)
            anch = {"l": "la", "c": "ma", "r": "ra"}[o["anchor"]]
            d.text((X(a[0]), X(a[1])), b, font=f, fill=col, anchor=anch)

    img.resize((int(sheet.w * dpi), int(sheet.h * dpi)),
               Image.LANCZOS).save(path, optimize=True)
    return int(sheet.w * dpi), int(sheet.h * dpi)


# ---------------------------------------------------------------------------
# 5.  PPTX backend
# ---------------------------------------------------------------------------


def render_pptx(sheet, path):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn

    def E(v):
        return Emu(int(round(v * 914400)))

    def C(h):
        return RGBColor.from_string(h)

    prs = Presentation()
    prs.slide_width = E(sheet.w)
    prs.slide_height = E(sheet.h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    def style(shape, o):
        ln = shape.line
        if o.get("noline"):
            ln.fill.background()
            return
        ln.color.rgb = C(o.get("color", INK))
        ln.width = Pt(o.get("pt", 1.0))
        el = ln._get_or_add_ln()
        for tag in ("a:prstDash", "a:headEnd", "a:tailEnd"):
            for e in el.findall(qn(tag)):
                el.remove(e)
        if o.get("dash"):
            el.append(el.makeelement(qn("a:prstDash"),
                                     {"val": DASHES[o["dash"]][2]}))
        for tag, on in (("a:headEnd", o.get("arrow_start")),
                        ("a:tailEnd", o.get("arrow_end"))):
            if on:
                el.append(el.makeelement(
                    qn(tag), {"type": "triangle", "w": "lg", "len": "lg"}))

    for kind, a, b, o in sheet.items:
        if kind == "line":
            sh = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      E(a[0]), E(a[1]), E(b[0]), E(b[1]))
            style(sh, o)
        elif kind == "circle":
            cx, cy, r = a[0], a[1], b
            sh = shapes.add_shape(MSO_SHAPE.OVAL, E(cx - r), E(cy - r),
                                  E(2 * r), E(2 * r))
            sh.shadow.inherit = False
            if o.get("fill"):
                sh.fill.solid()
                sh.fill.fore_color.rgb = C(o["fill"])
            else:
                sh.fill.background()
            style(sh, o)
        elif kind == "poly":
            pts = [(E(p[0]), E(p[1])) for p in a]
            ff = shapes.build_freeform(pts[0][0], pts[0][1])
            ff.add_line_segments(pts[1:], close=bool(o.get("close", True)))
            sh = ff.convert_to_shape()
            sh.shadow.inherit = False
            if o.get("fill"):
                sh.fill.solid()
                sh.fill.fore_color.rgb = C(o["fill"])
            else:
                sh.fill.background()
            style(sh, o)
        else:
            w = text_w(b, o["size"], o["bold"]) + 0.06
            lh = o["size"] * 1.35 / 72.0
            x = {"l": a[0], "c": a[0] - w / 2.0, "r": a[0] - w}[o["anchor"]]
            tb = shapes.add_textbox(E(x), E(a[1] - 0.02), E(w), E(lh + 0.04))
            tf = tb.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                           "r": PP_ALIGN.RIGHT}[o["anchor"]]
            run = p.add_run()
            run.text = b
            run.font.name = "Arial"
            run.font.size = Pt(o["size"])
            run.font.bold = o["bold"]
            run.font.color.rgb = C(o["color"])

    prs.save(path)


# ---------------------------------------------------------------------------

if __name__ == "__main__":
    dpi = 96
    if "--dpi" in sys.argv:
        dpi = float(sys.argv[sys.argv.index("--dpi") + 1])
    here = os.path.dirname(os.path.abspath(__file__))
    sheet = build()
    png = os.path.join(here, "dc_params_primer.png")
    w, h = render_png(sheet, png, dpi=dpi)
    print("PNG  %dx%d  ~%d image tokens  %d bytes"
          % (w, h, w * h / 750, os.path.getsize(png)))
    ppt = os.path.join(here, "dc_params_primer.pptx")
    render_pptx(sheet, ppt)
    print("PPTX %d bytes  (%d shapes)"
          % (os.path.getsize(ppt), len(sheet.items)))
